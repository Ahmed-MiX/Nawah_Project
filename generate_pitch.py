"""
Nawah Enterprise Pitch Deck Generator
Generates a professional, RTL-correct Arabic PowerPoint presentation.
Uses arabic_reshaper + bidi for perfect Arabic glyph rendering.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import arabic_reshaper
from bidi.algorithm import get_display


# ============================================================
# BRANDING CONSTANTS
# ============================================================
NAVY = RGBColor(10, 36, 99)
GREEN = RGBColor(52, 168, 83)
GOLD = RGBColor(249, 171, 0)
PURPLE = RGBColor(128, 0, 128)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 245)
DARK_GRAY = RGBColor(80, 80, 80)


# ============================================================
# ARABIC TEXT HELPER — MANDATORY for every Arabic string
# ============================================================
def fix_arabic(text):
    """Reshape and reorder Arabic text for correct RTL rendering in PPTX."""
    reshaped = arabic_reshaper.reshape(text)
    return get_display(reshaped)


# ============================================================
# SLIDE HELPERS
# ============================================================
def set_slide_bg(slide, color):
    """Set the background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, color, left, top, width, height):
    """Add a colored accent rectangle bar to a slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # No border
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, font_color=NAVY, bold=False, alignment=PP_ALIGN.RIGHT):
    """Add a right-aligned Arabic text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = fix_arabic(text)
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=16, font_color=DARK_GRAY, bullet_char="◂"):
    """Add a right-aligned bulleted list of Arabic items."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = fix_arabic(f"{bullet_char} {item}")
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.RIGHT
        p.space_after = Pt(8)

    return txBox


def add_card(slide, title, body, left, top, width, height,
             accent_color=GREEN, title_size=16, body_size=13):
    """Add a card-style box with accent top bar, title, and body text."""
    # Card background
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(1)
    card.shadow.inherit = False

    # Accent top bar
    bar_height = Pt(5)
    add_accent_bar(slide, accent_color, left, top, width, bar_height)

    # Title text
    add_text_box(slide, title,
                 left + Inches(0.15), top + Pt(14),
                 width - Inches(0.3), Inches(0.4),
                 font_size=title_size, font_color=NAVY, bold=True)

    # Body text
    add_text_box(slide, body,
                 left + Inches(0.15), top + Inches(0.55),
                 width - Inches(0.3), height - Inches(0.7),
                 font_size=body_size, font_color=DARK_GRAY)


# ============================================================
# SLIDE BUILDERS
# ============================================================
def build_slide_1_title(prs):
    """Slide 1: Title — The Grand Opening."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, WHITE)

    # Top gold accent bar
    add_accent_bar(slide, GOLD, Inches(0), Inches(0), Inches(13.33), Pt(6))

    # Bottom navy accent bar
    add_accent_bar(slide, NAVY, Inches(0), Inches(7.2), Inches(13.33), Pt(12))

    # Left green vertical accent
    add_accent_bar(slide, GREEN, Inches(0.4), Inches(1.5), Pt(6), Inches(4))

    # Main title
    add_text_box(slide, "نَوَاة",
                 Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
                 font_size=60, font_color=NAVY, bold=True)

    # Subtitle
    add_text_box(slide, "مركز القيادة والسيطرة لإدارة الكيانات ذاتياً",
                 Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
                 font_size=28, font_color=PURPLE, bold=True)

    # Tagline
    add_text_box(slide, "محرك تشغيلي ذكي يحوّل أي مؤسسة إلى كيان ذاتي الإدارة",
                 Inches(1.5), Inches(4.0), Inches(10), Inches(0.6),
                 font_size=18, font_color=DARK_GRAY)

    # Decorative accent dots
    add_accent_bar(slide, GREEN, Inches(11.5), Inches(2.0), Pt(12), Pt(12))
    add_accent_bar(slide, GOLD, Inches(11.8), Inches(2.0), Pt(12), Pt(12))
    add_accent_bar(slide, PURPLE, Inches(12.1), Inches(2.0), Pt(12), Pt(12))


def build_slide_2_problem(prs):
    """Slide 2: The Problem — Why Nawah Exists."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Top accent bar
    add_accent_bar(slide, NAVY, Inches(0), Inches(0), Inches(13.33), Pt(4))

    # Section label
    add_text_box(slide, "المشكلة",
                 Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
                 font_size=14, font_color=GOLD, bold=True)

    # Title
    add_text_box(slide, "الشلل التشغيلي في المؤسسات",
                 Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 font_size=36, font_color=NAVY, bold=True)

    # Problem cards — 3 columns
    card_w = Inches(3.5)
    card_h = Inches(3.5)
    card_y = Inches(2.2)

    add_card(slide,
             "الموظف كجهاز توجيه",
             "كل موظف يعمل كراوتر بشري: يستقبل المعلومات، يعالجها يدوياً، ويعيد توجيهها. عند غيابه ينهار المسار بالكامل.",
             Inches(0.8), card_y, card_w, card_h,
             accent_color=GOLD, title_size=18, body_size=14)

    add_card(slide,
             "تشتت السياق المؤسسي",
             "المعرفة المؤسسية مبعثرة بين إيميلات وملفات ومحادثات. لا يوجد عقل مركزي يربط السياقات ببعضها.",
             Inches(4.8), card_y, card_w, card_h,
             accent_color=PURPLE, title_size=18, body_size=14)

    add_card(slide,
             "الاعتماد على الذاكرة البشرية",
             "القرارات تعتمد على ما يتذكره الأفراد. عند انتقال موظف تضيع سنوات من السياق التشغيلي.",
             Inches(8.8), card_y, card_w, card_h,
             accent_color=GREEN, title_size=18, body_size=14)

    # Bottom separator
    add_accent_bar(slide, LIGHT_GRAY, Inches(0.8), Inches(6.5), Inches(11.5), Pt(1))

    # Impact statement
    add_text_box(slide, "النتيجة: مؤسسات تعمل بكفاءة 30% فقط من طاقتها الحقيقية",
                 Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
                 font_size=16, font_color=PURPLE, bold=True)


def build_slide_3_solution(prs):
    """Slide 3: The Solution — What Nawah Does."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Top accent bar
    add_accent_bar(slide, GREEN, Inches(0), Inches(0), Inches(13.33), Pt(4))

    # Section label
    add_text_box(slide, "الحل",
                 Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
                 font_size=14, font_color=GREEN, bold=True)

    # Title
    add_text_box(slide, "نَوَاة: المحرك التشغيلي المستقل",
                 Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 font_size=36, font_color=NAVY, bold=True)

    # Left column: Core capabilities
    capabilities = [
        "استشعار شامل: البريد، الملفات، الأوامر الصوتية والنصية",
        "تحليل ذكي: تصنيف المهام وتحديد الوكلاء تلقائياً",
        "تنفيذ مستقل: سرب وكلاء متخصصين ينفذون بلا تدخل بشري",
        "رؤية حاسوبية: تحليل الصور والمستندات الممسوحة ضوئياً",
        "ذاكرة مؤسسية: سياق تراكمي لا يُفقد أبداً",
    ]

    add_bullet_list(slide, capabilities,
                    Inches(0.8), Inches(2.2), Inches(5.5), Inches(4),
                    font_size=16, font_color=DARK_GRAY, bullet_char="◂")

    # Right column: Architecture card
    add_card(slide,
             "البنية التحتية",
             "Gemini 2.5 Flash + LangChain + Streamlit + Watchdog + IMAP Radar",
             Inches(7), Inches(2.2), Inches(5), Inches(1.2),
             accent_color=NAVY, title_size=16, body_size=13)

    add_card(slide,
             "التدفق التشغيلي",
             "أمر → تحليل L1 → تصنيف المهمة → توزيع على الوكلاء → تنفيذ → تقرير JSON",
             Inches(7), Inches(3.8), Inches(5), Inches(1.2),
             accent_color=GREEN, title_size=16, body_size=13)

    add_card(slide,
             "القنوات المدعومة",
             "واجهة ويب تفاعلية • رادار بريد إلكتروني • حارس مجلدات رقابي",
             Inches(7), Inches(5.4), Inches(5), Inches(1.2),
             accent_color=GOLD, title_size=16, body_size=13)


def build_slide_4_resilience(prs):
    """Slide 4: System Resilience — Anti-Fragile Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Top accent bar
    add_accent_bar(slide, PURPLE, Inches(0), Inches(0), Inches(13.33), Pt(4))

    # Section label
    add_text_box(slide, "المتانة",
                 Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
                 font_size=14, font_color=PURPLE, bold=True)

    # Title
    add_text_box(slide, "هندسة مضادة للانهيار",
                 Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 font_size=36, font_color=NAVY, bold=True)

    # Resilience features — 4 cards in 2x2 grid
    card_w = Inches(5.3)
    card_h = Inches(1.8)

    add_card(slide,
             "صفر تأخير في الاستجابة",
             "المسح الاستهلالي يعالج كل الملفات المتأخرة عند إعادة التشغيل. لا يضيع أي ملف حتى لو سقط النظام.",
             Inches(0.8), Inches(2.2), card_w, card_h,
             accent_color=GREEN)

    add_card(slide,
             "مضاد للاختناق (Anti-Crash)",
             "عزل كامل لمعالجة كل إيميل. فشل واحد لا يوقف البقية. تبريد 10 ثوانٍ بين كل عملية.",
             Inches(6.8), Inches(2.2), card_w, card_h,
             accent_color=GOLD)

    add_card(slide,
             "جدار حماية ذكي",
             "فلترة الملفات حسب الصيغة (9 أنواع مدعومة) والحجم (20MB كحد أقصى). رفض فوري للملفات المشبوهة.",
             Inches(0.8), Inches(4.5), card_w, card_h,
             accent_color=PURPLE)

    add_card(slide,
             "عزل مطلق للبيانات",
             "UUID فريد لكل ملف مؤقت. صفر تصادم حتى مع مرفقات متطابقة الأسماء من إيميلات متزامنة.",
             Inches(6.8), Inches(4.5), card_w, card_h,
             accent_color=NAVY)

    # Bottom confidence bar
    add_accent_bar(slide, GREEN, Inches(0.8), Inches(6.8), Inches(11.5), Pt(3))
    add_text_box(slide, "معدل التوافر المستهدف: 99.9% — تصميم مبني على مبدأ الفشل الآمن",
                 Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
                 font_size=14, font_color=GREEN, bold=True)


def build_slide_5_scalability(prs):
    """Slide 5: Scalability — Cross-Industry Adaptability."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Top accent bar
    add_accent_bar(slide, GOLD, Inches(0), Inches(0), Inches(13.33), Pt(4))

    # Section label
    add_text_box(slide, "قابلية التوسع",
                 Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
                 font_size=14, font_color=GOLD, bold=True)

    # Title
    add_text_box(slide, "منظومة محايدة تتكيف مع أي قطاع",
                 Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 font_size=36, font_color=NAVY, bold=True)

    # Industry cards — 3 columns
    card_w = Inches(3.5)
    card_h = Inches(3.2)
    card_y = Inches(2.2)

    add_card(slide,
             "القطاع المالي والمصرفي",
             "أتمتة طلبات القروض، تحليل المخاطر، توجيه المعاملات، ومراقبة الامتثال التنظيمي بدون تدخل بشري.",
             Inches(0.8), card_y, card_w, card_h,
             accent_color=NAVY, title_size=18, body_size=14)

    add_card(slide,
             "القطاع الصحي",
             "جدولة المواعيد، تصنيف التقارير الطبية، تحليل الأشعة بالذكاء الاصطناعي، وإدارة ملفات المرضى.",
             Inches(4.8), card_y, card_w, card_h,
             accent_color=GREEN, title_size=18, body_size=14)

    add_card(slide,
             "القطاع الحكومي",
             "معالجة المعاملات الإلكترونية، التوجيه الآلي للطلبات، تحليل الوثائق الرسمية، وإعداد التقارير.",
             Inches(8.8), card_y, card_w, card_h,
             accent_color=PURPLE, title_size=18, body_size=14)

    # Scalability statement
    add_accent_bar(slide, LIGHT_GRAY, Inches(0.8), Inches(6.0), Inches(11.5), Pt(1))

    add_text_box(slide, "نَوَاة لا تحتاج إعادة برمجة — فقط إعادة تهيئة السياق للقطاع الجديد",
                 Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
                 font_size=16, font_color=NAVY, bold=True)


def build_slide_6_roadmap(prs):
    """Slide 6: Roadmap — The Future Vision."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Top accent bar
    add_accent_bar(slide, NAVY, Inches(0), Inches(0), Inches(13.33), Pt(4))

    # Section label
    add_text_box(slide, "خارطة الطريق",
                 Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
                 font_size=14, font_color=NAVY, bold=True)

    # Title
    add_text_box(slide, "الآلة تنفذ والقائد يعتمد",
                 Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 font_size=36, font_color=NAVY, bold=True)

    # Timeline phases — horizontal flow
    phase_w = Inches(3.5)
    phase_h = Inches(3.0)
    phase_y = Inches(2.3)

    add_card(slide,
             "المرحلة 1: الاندماج العضوي",
             "بناء النواة الأساسية: واجهة القيادة، محلل المهام L1، رادار البريد، الحارس الرقابي، والرؤية الحاسوبية. — مكتمل ✓",
             Inches(0.8), phase_y, phase_w, phase_h,
             accent_color=GREEN, title_size=16, body_size=13)

    add_card(slide,
             "المرحلة 2: الذكاء الاستباقي",
             "وكلاء متخصصون (مالي، قانوني، تقني)، ذاكرة سياقية طويلة المدى، تعلم من الأنماط التشغيلية المتكررة.",
             Inches(4.8), phase_y, phase_w, phase_h,
             accent_color=GOLD, title_size=16, body_size=13)

    add_card(slide,
             "المرحلة 3: الحكم الذاتي الكامل",
             "اتخاذ قرارات تشغيلية مستقلة، تكامل مع أنظمة ERP وCRM، لوحة قيادة تنفيذية للإدارة العليا.",
             Inches(8.8), phase_y, phase_w, phase_h,
             accent_color=PURPLE, title_size=16, body_size=13)

    # Connecting arrows between phases
    add_accent_bar(slide, GOLD, Inches(4.4), Inches(3.7), Inches(0.3), Pt(4))
    add_accent_bar(slide, GOLD, Inches(8.4), Inches(3.7), Inches(0.3), Pt(4))

    # Bottom closing
    add_accent_bar(slide, NAVY, Inches(0), Inches(7.0), Inches(13.33), Pt(8))
    add_text_box(slide, "نَوَاة — حيث تتحول المؤسسات من التشغيل اليدوي إلى الحكم الذاتي",
                 Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
                 font_size=18, font_color=NAVY, bold=True)


# ============================================================
# MAIN EXECUTION
# ============================================================
def generate_presentation():
    """Generate the complete Nawah Enterprise Pitch Deck."""
    prs = Presentation()

    # Set widescreen 16:9 dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Build all 6 slides in order
    print("📊 جاري إنشاء العرض التقديمي...")

    build_slide_1_title(prs)
    print("  ✅ الشريحة 1: العنوان الرئيسي")

    build_slide_2_problem(prs)
    print("  ✅ الشريحة 2: المشكلة")

    build_slide_3_solution(prs)
    print("  ✅ الشريحة 3: الحل")

    build_slide_4_resilience(prs)
    print("  ✅ الشريحة 4: المتانة")

    build_slide_5_scalability(prs)
    print("  ✅ الشريحة 5: قابلية التوسع")

    build_slide_6_roadmap(prs)
    print("  ✅ الشريحة 6: خارطة الطريق")

    # Save the file
    output_path = "Nawah_Enterprise_Pitch.pptx"
    prs.save(output_path)
    print(f"\n🎯 تم إنشاء العرض التقديمي بنجاح: {output_path}")

    return output_path


if __name__ == "__main__":
    try:
        generate_presentation()
    except Exception as e:
        print(f"❌ خطأ في إنشاء العرض: {e}")
        raise
