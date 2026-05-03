from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import arabic_reshaper
from bidi.algorithm import get_display
import os

NAVY = RGBColor(10, 36, 99)
GREEN = RGBColor(52, 168, 83)
GOLD = RGBColor(249, 171, 0)
PURPLE = RGBColor(128, 0, 128)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(245, 247, 250)
DGRAY = RGBColor(70, 70, 80)
LGRAY = RGBColor(225, 228, 232)
FONT = "Segoe UI"
LOGO = os.path.join(os.path.dirname(__file__), "شعار مشروع نواة المبدئي.jpeg")

def fa(text):
    return get_display(arabic_reshaper.reshape(text))

def set_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG

def bar(slide, color, l, t, w, h):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def logo_corner(slide):
    try:
        slide.shapes.add_picture(LOGO, Inches(11.3), Inches(0.2), height=Inches(0.7))
    except Exception:
        pass

def logo_big(slide):
    try:
        slide.shapes.add_picture(LOGO, Inches(4.8), Inches(0.6), height=Inches(2.8))
    except Exception:
        pass

def tbox(slide, text, l, t, w, h, sz=18, clr=NAVY, bold=False, align=PP_ALIGN.RIGHT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = fa(text)
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    return tb

def bullets(slide, items, l, t, w, h, sz=15, clr=DGRAY, ch="◂"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = fa(f"{ch} {item}")
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.name = FONT
        p.alignment = PP_ALIGN.RIGHT
        p.space_after = Pt(6)
    return tb

def card(slide, title, body, l, t, w, h, ac=GREEN, tsz=15, bsz=12):
    c = slide.shapes.add_shape(1, l, t, w, h)
    c.fill.solid()
    c.fill.fore_color.rgb = WHITE
    c.line.color.rgb = LGRAY
    c.line.width = Pt(1)
    bar(slide, ac, l, t, w, Pt(5))
    tbox(slide, title, l+Inches(0.12), t+Pt(12), w-Inches(0.24), Inches(0.35), sz=tsz, clr=NAVY, bold=True)
    tbox(slide, body, l+Inches(0.12), t+Inches(0.5), w-Inches(0.24), h-Inches(0.6), sz=bsz, clr=DGRAY)
