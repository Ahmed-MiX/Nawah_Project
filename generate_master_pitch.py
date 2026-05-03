from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pitch_helpers import *

def s1(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    bar(sl, GOLD, Inches(0), Inches(0), Inches(13.33), Pt(6))
    bar(sl, NAVY, Inches(0), Inches(7.2), Inches(13.33), Pt(10))
    bar(sl, GREEN, Inches(0.5), Inches(3.6), Pt(5), Inches(2.5))
    logo_big(sl)
    tbox(sl, "نَوَاة", Inches(1), Inches(3.6), Inches(11), Inches(0.9), sz=54, clr=NAVY, bold=True, align=PP_ALIGN.CENTER)
    tbox(sl, "مركز القيادة والسيطرة التشغيلي", Inches(1), Inches(4.5), Inches(11), Inches(0.6), sz=26, clr=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    tbox(sl, "أتمتة الكيانات عبر الأنظمة متعددة الوكلاء", Inches(1), Inches(5.2), Inches(11), Inches(0.5), sz=18, clr=DGRAY, align=PP_ALIGN.CENTER)
    bar(sl, GREEN, Inches(5.5), Inches(6.0), Inches(2.3), Pt(3))

def s2(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl); logo_corner(sl)
    bar(sl, GOLD, Inches(0), Inches(0), Inches(13.33), Pt(4))
    tbox(sl, "الأزمة الخفية", Inches(0.8), Inches(0.3), Inches(3), Inches(0.4), sz=13, clr=GOLD, bold=True)
    tbox(sl, "الشلل التشغيلي: لماذا تفشل المؤسسات من الداخل؟", Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), sz=32, clr=NAVY, bold=True)
    cw, ch, cy = Inches(3.5), Inches(3.8), Inches(2.0)
    card(sl, "الموظف كجهاز توجيه", "كل موظف يعمل كراوتر بشري: يستقبل المعلومات، يعالجها يدوياً، ويعيد توجيهها لجهة أخرى. عند غيابه ينهار المسار بالكامل وتتوقف سلاسل القرار. المؤسسة لا تعمل — بل تنتظر موظفين ليوجّهوا البيانات.", Inches(0.8), cy, cw, ch, ac=GOLD, tsz=17, bsz=13)
    card(sl, "تشتت السياق المؤسسي", "المعرفة التشغيلية مبعثرة بين إيميلات ومحادثات وملفات متفرقة. لا يوجد عقل مركزي يربط السياقات. كل قسم يعمل في جزيرة معزولة. عند انتقال موظف تضيع سنوات من الخبرة التراكمية بلا رجعة.", Inches(4.8), cy, cw, ch, ac=PURPLE, tsz=17, bsz=13)
    card(sl, "فشل الأتمتة العمياء", "أدوات الأتمتة التقليدية تنفذ أوامر ثابتة بلا فهم. لا تستوعب السياق ولا تتكيف مع المتغيرات. تحتاج برمجة مسبقة لكل سيناريو — وعند ظهور حالة جديدة تتوقف تماماً وتحتاج تدخل بشري.", Inches(8.8), cy, cw, ch, ac=GREEN, tsz=17, bsz=13)
    bar(sl, LGRAY, Inches(0.8), Inches(6.3), Inches(11.5), Pt(1))
    tbox(sl, "النتيجة: مؤسسات تعمل بكفاءة 30% فقط من طاقتها الحقيقية", Inches(0.8), Inches(6.5), Inches(11), Inches(0.4), sz=15, clr=PURPLE, bold=True)

def s3(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl); logo_corner(sl)
    bar(sl, GREEN, Inches(0), Inches(0), Inches(13.33), Pt(4))
    tbox(sl, "الحل", Inches(0.8), Inches(0.3), Inches(3), Inches(0.4), sz=13, clr=GREEN, bold=True)
    tbox(sl, "نَوَاة: من أدوات جامدة إلى قوى عاملة ذكية مستقلة", Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), sz=30, clr=NAVY, bold=True)
    tbox(sl, "نَوَاة ليست أداة أتمتة — بل محرك تشغيلي ذكي يحوّل أي مؤسسة إلى كيان ذاتي الإدارة. تستقبل المهام من كل القنوات، تفهم السياق، تخطط، وتنفذ عبر سرب من الوكلاء المتخصصين.", Inches(0.8), Inches(1.8), Inches(11), Inches(0.8), sz=16, clr=DGRAY)
    bullets(sl, [
        "استشعار شامل: بريد إلكتروني، ملفات، صور، أوامر نصية — كلها تتحول لسياق موحد",
        "حقيبة السياق الموحدة: كل المدخلات تُدمج في وعاء واحد قبل التحليل",
        "تحليل ذكي L1: تصنيف المهام وتحديد الوكلاء المطلوبين تلقائياً بالذكاء الاصطناعي",
        "تنفيذ مستقل: سرب وكلاء متخصصين ينفذون بلا تدخل بشري",
        "ذاكرة مؤسسية تراكمية: السياق لا يُفقد أبداً حتى مع تغيّر الموظفين",
    ], Inches(0.8), Inches(2.8), Inches(5.8), Inches(4), sz=15)
    card(sl, "البنية التحتية", "Gemini 2.5 Flash + LangChain + Streamlit + Watchdog + IMAP", Inches(7.2), Inches(2.8), Inches(5), Inches(1.1), ac=NAVY, tsz=15, bsz=12)
    card(sl, "التدفق التشغيلي", "أمر → تحليل L1 → تصنيف → توزيع على الوكلاء → تنفيذ → تقرير JSON", Inches(7.2), Inches(4.2), Inches(5), Inches(1.1), ac=GREEN, tsz=15, bsz=12)
    card(sl, "القنوات المدعومة", "واجهة ويب تفاعلية • رادار بريد إلكتروني • حارس مجلدات رقابي • رؤية حاسوبية", Inches(7.2), Inches(5.6), Inches(5), Inches(1.1), ac=GOLD, tsz=15, bsz=12)

def s4(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl); logo_corner(sl)
    bar(sl, NAVY, Inches(0), Inches(0), Inches(13.33), Pt(4))
    tbox(sl, "الطبقة الأولى", Inches(0.8), Inches(0.3), Inches(3), Inches(0.4), sz=13, clr=NAVY, bold=True)
    tbox(sl, "الاستقبال الذكي: رادار الاستشعار متعدد القنوات", Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), sz=30, clr=NAVY, bold=True)
    tbox(sl, "الطبقة الأولى تعمل كجهاز عصبي حسّي — تلتقط كل مدخل من أي قناة وتحوّله إلى نص مفهوم قبل إرساله للعقل المحلل. لا يوجد مدخل يضيع.", Inches(0.8), Inches(1.8), Inches(11), Inches(0.7), sz=15, clr=DGRAY)
    cw, ch2, cy = Inches(3.5), Inches(3.2), Inches(2.8)
    card(sl, "رادار البريد الإلكتروني", "مراقبة مستمرة للبريد كل 15 ثانية عبر IMAP. استخراج النص والمرفقات تلقائياً. فك تشفير العناوين والمحتوى العربي بالكامل. دمج جسم الرسالة مع كل المرفقات في استعلام واحد موحد.", Inches(0.8), cy, cw, ch2, ac=GOLD, tsz=16, bsz=13)
    card(sl, "الرؤية الحاسوبية", "تحليل الصور والمستندات الممسوحة ضوئياً عبر Gemini 2.5 Flash Vision. استخراج النصوص والأشكال والرسوم البيانية. تحليل عميق وشامل لمحتوى كل صورة مع وصف تفصيلي لنية الملف الأصلية.", Inches(4.8), cy, cw, ch2, ac=PURPLE, tsz=16, bsz=13)
    card(sl, "الحارس الرقابي للمجلدات", "مراقبة فورية لمجلد nawah_inbox عبر Watchdog. مسح استهلالي عند التشغيل لالتقاط الملفات المتأخرة. دعم 9 صيغ: PDF, TXT, CSV, DOCX, XLSX, PNG, JPG, JPEG. نقل الملفات المعالجة للأرشيف.", Inches(8.8), cy, cw, ch2, ac=GREEN, tsz=16, bsz=13)

def s5(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl); logo_corner(sl)
    bar(sl, PURPLE, Inches(0), Inches(0), Inches(13.33), Pt(4))
    tbox(sl, "الطبقة الثانية", Inches(0.8), Inches(0.3), Inches(3), Inches(0.4), sz=13, clr=PURPLE, bold=True)
    tbox(sl, "التخطيط والتفويض: العقل المحلل L1", Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), sz=30, clr=NAVY, bold=True)
    tbox(sl, "بعد استقبال المدخلات الموحدة، يقوم المحلل L1 بتفكيك المهمة وتحديد مسار التنفيذ الاستراتيجي.", Inches(0.8), Inches(1.7), Inches(11), Inches(0.6), sz=15, clr=DGRAY)
    card(sl, "تحليل النية", "فهم عميق لمقصد المستخدم من النص والمرفقات. تلخيص النية في جملة عربية واضحة تحدد بدقة ما يريده المستخدم.", Inches(0.8), Inches(2.6), Inches(5.3), Inches(1.6), ac=NAVY, tsz=16, bsz=13)
    card(sl, "تحديد الوكلاء المطلوبين", "تصنيف ذكي للأدوار المطلوبة: محلل نصوص، خبير أتمتة، مهندس نظم، محلل صور — كل مهمة تحصل على فريقها المناسب.", Inches(6.8), Inches(2.6), Inches(5.3), Inches(1.6), ac=GREEN, tsz=16, bsz=13)
    card(sl, "تقييم التعقيد المعاير", "مقياس ثلاثي صارم: Low للمهام البسيطة، Medium للتقارير والتحليلات، High فقط للقرارات الاستراتيجية والمالية الكبرى. معاير لمنع التضخم.", Inches(0.8), Inches(4.6), Inches(5.3), Inches(1.6), ac=GOLD, tsz=16, bsz=13)
    card(sl, "المخرج الموحد: JSON", "كل تحليل ينتج ملف JSON منظم يحتوي: intent, agents_needed, complexity. مخرج قابل للقراءة الآلية والبشرية على حد سواء.", Inches(6.8), Inches(4.6), Inches(5.3), Inches(1.6), ac=PURPLE, tsz=16, bsz=13)

def s6(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl); logo_corner(sl)
    bar(sl, GREEN, Inches(0), Inches(0), Inches(13.33), Pt(4))
    tbox(sl, "الطبقة الثالثة", Inches(0.8), Inches(0.3), Inches(3), Inches(0.4), sz=13, clr=GREEN, bold=True)
    tbox(sl, "العصب التنفيذي: سرب الوكلاء الذكي", Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), sz=30, clr=NAVY, bold=True)
    tbox(sl, "بعد انتهاء التخطيط، يتولى سرب من الوكلاء المتخصصين تنفيذ المهام بشكل متوازٍ ومستقل. كل وكيل يملك تخصصاً محدداً ويتواصل مع بقية السرب لضمان التنسيق الكامل.", Inches(0.8), Inches(1.7), Inches(11), Inches(0.7), sz=15, clr=DGRAY)
    bullets(sl, [
        "وكيل تحليل النصوص: استخراج المعلومات الجوهرية من المستندات والتقارير",
        "وكيل الرؤية الحاسوبية: تحليل الصور والمخططات والجداول المصورة",
        "وكيل التخطيط الاستراتيجي: بناء خطط التنفيذ وتوزيع المهام الفرعية",
        "وكيل التواصل: صياغة الردود والتقارير باللغة العربية الاحترافية",
        "وكيل المراقبة: تتبع حالة التنفيذ وإبلاغ مركز القيادة بالنتائج",
    ], Inches(0.8), Inches(2.8), Inches(6), Inches(3.5), sz=14)
    card(sl, "التنسيق بين الوكلاء", "كل وكيل يعمل بشكل مستقل لكنه يشارك النتائج مع بقية السرب. لا يوجد اختناق مركزي — التنفيذ متوازٍ وسريع.", Inches(7.5), Inches(2.8), Inches(4.8), Inches(1.5), ac=NAVY)
    card(sl, "سرعة التنفيذ", "من استقبال الأمر إلى تسليم النتيجة في ثوانٍ معدودة. المهام التي تستغرق ساعات بشرياً تُنجز آلياً في لحظات.", Inches(7.5), Inches(4.6), Inches(4.8), Inches(1.5), ac=GOLD)

def s7(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl); logo_corner(sl)
    bar(sl, GOLD, Inches(0), Inches(0), Inches(13.33), Pt(4))
    tbox(sl, "القرار السيادي", Inches(0.8), Inches(0.3), Inches(3), Inches(0.4), sz=13, clr=GOLD, bold=True)
    tbox(sl, "التوازن بين الأتمتة المطلقة والسيطرة البشرية", Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), sz=30, clr=NAVY, bold=True)
    tbox(sl, "نَوَاة لا تستبدل القائد — بل تحرره. المهام الروتينية تُنفَّذ تلقائياً، والقرارات الاستراتيجية تُعرض على لوحة القيادة للاعتماد البشري.", Inches(0.8), Inches(1.7), Inches(11), Inches(0.7), sz=15, clr=DGRAY)
    card(sl, "الأتمتة الكاملة للمهام الروتينية", "المهام المصنفة Low و Medium تُنفَّذ فوراً بدون تدخل بشري. تحليل الملفات، الرد على الاستفسارات، تصنيف البريد — كلها تعمل في الخلفية بصمت تام.", Inches(0.8), Inches(2.8), Inches(5.3), Inches(2), ac=GREEN, tsz=16, bsz=13)
    card(sl, "بوابة الاعتماد الاستراتيجي", "المهام المصنفة High تتوقف عند بوابة الاعتماد. يراها القائد على لوحة القيادة مع التحليل الكامل والتوصيات. يقرر: تنفيذ، تعديل، أو رفض.", Inches(6.8), Inches(2.8), Inches(5.3), Inches(2), ac=GOLD, tsz=16, bsz=13)
    card(sl, "لوحة القيادة التنفيذية", "واجهة Streamlit تفاعلية تعرض: حالة الوكلاء، المهام المؤتمتة في الخلفية، نتائج التحليل بصيغة JSON، وسجل كامل لكل العمليات. رؤية شاملة بنقرة واحدة.", Inches(0.8), Inches(5.2), Inches(11.3), Inches(1.6), ac=NAVY, tsz=16, bsz=13)

def s8(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl); logo_corner(sl)
    bar(sl, PURPLE, Inches(0), Inches(0), Inches(13.33), Pt(4))
    tbox(sl, "المتانة المؤسسية", Inches(0.8), Inches(0.3), Inches(3), Inches(0.4), sz=13, clr=PURPLE, bold=True)
    tbox(sl, "هندسة مضادة للانهيار: صفر توقف، صفر فقدان", Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), sz=30, clr=NAVY, bold=True)
    cw, ch2 = Inches(5.3), Inches(1.7)
    card(sl, "المسح الاستهلالي عند التشغيل", "عند إعادة تشغيل النظام، يفحص الحارس الرقابي كل الملفات المتأخرة في صندوق الوارد ويعالجها فوراً. لا يضيع أي ملف حتى لو سقط النظام.", Inches(0.8), Inches(2.0), cw, ch2, ac=GREEN)
    card(sl, "العزل الكامل لكل عملية (Anti-Crash)", "كل إيميل يُعالج في كبسولة معزولة. فشل واحد لا يوقف البقية. تبريد 10 ثوانٍ بين كل عملية لحماية من حدود API.", Inches(6.8), Inches(2.0), cw, ch2, ac=GOLD)
    card(sl, "صفر تصادم بيانات (UUID)", "كل ملف مؤقت يحصل على معرّف UUID فريد. حتى لو وصلت مرفقات بنفس الاسم من إيميلات متزامنة — لا يوجد أي تصادم أو فقدان.", Inches(0.8), Inches(4.1), cw, ch2, ac=PURPLE)
    card(sl, "جدار حماية ذكي صامت", "فلترة تلقائية حسب الصيغة (9 أنواع) والحجم (20MB). رفض فوري للملفات المشبوهة أو الضخمة مع تسجيل سبب الرفض بالعربية.", Inches(6.8), Inches(4.1), cw, ch2, ac=NAVY)
    bar(sl, GREEN, Inches(0.8), Inches(6.3), Inches(11.5), Pt(3))
    tbox(sl, "معدل التوافر المستهدف: 99.9% — تصميم مبني على مبدأ الفشل الآمن", Inches(0.8), Inches(6.5), Inches(11), Inches(0.4), sz=14, clr=GREEN, bold=True)

def s9(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl); logo_corner(sl)
    bar(sl, NAVY, Inches(0), Inches(0), Inches(13.33), Pt(4))
    tbox(sl, "المستقبل", Inches(0.8), Inches(0.3), Inches(3), Inches(0.4), sz=13, clr=NAVY, bold=True)
    tbox(sl, "بنية محايدة تتكيف مع أي قطاع — الآلة تنفذ والقائد يعتمد", Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), sz=28, clr=NAVY, bold=True)
    pw, ph, py = Inches(3.5), Inches(2.5), Inches(1.9)
    card(sl, "المرحلة 1: الاندماج العضوي ✓", "بناء النواة: واجهة القيادة، محلل L1، رادار البريد، الحارس الرقابي، الرؤية الحاسوبية، جدار الحماية. مكتمل.", Inches(0.8), py, pw, ph, ac=GREEN, tsz=15, bsz=12)
    card(sl, "المرحلة 2: الذكاء الاستباقي", "وكلاء متخصصون (مالي، قانوني، تقني). ذاكرة سياقية طويلة المدى. تعلم من الأنماط التشغيلية المتكررة واقتراح حلول استباقية.", Inches(4.8), py, pw, ph, ac=GOLD, tsz=15, bsz=12)
    card(sl, "المرحلة 3: الحكم الذاتي", "اتخاذ قرارات مستقلة. تكامل مع ERP و CRM. لوحة قيادة تنفيذية للإدارة العليا. أتمتة فائقة شاملة.", Inches(8.8), py, pw, ph, ac=PURPLE, tsz=15, bsz=12)
    bar(sl, GOLD, Inches(4.4), Inches(3.0), Inches(0.3), Pt(4))
    bar(sl, GOLD, Inches(8.4), Inches(3.0), Inches(0.3), Pt(4))
    tbox(sl, "القطاعات المستهدفة", Inches(0.8), Inches(4.8), Inches(11), Inches(0.4), sz=18, clr=NAVY, bold=True)
    cw2, ch3, cy2 = Inches(3.5), Inches(1.6), Inches(5.3)
    card(sl, "القطاع المالي والمصرفي", "أتمتة طلبات القروض، تحليل المخاطر، مراقبة الامتثال.", Inches(0.8), cy2, cw2, ch3, ac=NAVY, tsz=14, bsz=12)
    card(sl, "القطاع الصحي", "جدولة المواعيد، تصنيف التقارير الطبية، تحليل الأشعة.", Inches(4.8), cy2, cw2, ch3, ac=GREEN, tsz=14, bsz=12)
    card(sl, "القطاع الحكومي", "معالجة المعاملات الإلكترونية، التوجيه الآلي، إعداد التقارير.", Inches(8.8), cy2, cw2, ch3, ac=PURPLE, tsz=14, bsz=12)

def generate():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    for i, fn in enumerate([s1,s2,s3,s4,s5,s6,s7,s8,s9], 1):
        fn(prs)
        print(f"  [OK] Slide {i}")
    out = "Nawah_Master_Pitch.pptx"
    prs.save(out)
    print(f"\n  [DONE] {out}")

if __name__ == "__main__":
    generate()
